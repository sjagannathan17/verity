"""Generate the Verity slide deck (slides/Verity.pptx).

Enterprise, restrained, Cotiviti/Verity palette. 12 frames mirroring the
big-picture -> zoom-in -> pull-back narrative. Frames 2, 9, 10 carry a small
"[I will refine]" marker per the author's request.

Run:  ../.venv-tools/bin/python tools/build_slides.py   (from verity/)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "slides" / "Verity.pptx"

PURPLE = RGBColor(0x5A, 0x2D, 0x8C)
PURPLE_LT = RGBColor(0xF4, 0xF0, 0xFA)
CORAL = RGBColor(0xE8, 0x45, 0x3C)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    s = slide.shapes.add_shape(1, 0, 0, SW, SH)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shp = slide.shapes.add_shape(5 if radius else 1, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for item in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = item.get("align", align)
        p.space_after = Pt(item.get("space_after", 6))
        p.space_before = Pt(item.get("space_before", 0))
        if item.get("line_spacing"):
            p.line_spacing = item["line_spacing"]
        txt = item["text"]
        if item.get("bullet"):
            txt = "•   " + txt
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(item.get("size", 18))
        r.font.bold = item.get("bold", False)
        r.font.italic = item.get("italic", False)
        r.font.name = FONT
        r.font.color.rgb = item.get("color", SLATE)
    return tb


def kicker(slide, label):
    bar = box(slide, Inches(0.6), Inches(0.62), Inches(0.16), Inches(0.42), fill=CORAL, radius=False)
    bar.line.fill.background()
    text(slide, Inches(0.85), Inches(0.55), Inches(10), Inches(0.6),
         [{"text": label, "size": 13, "bold": True, "color": PURPLE}])


def heading(slide, title, sub=None):
    text(slide, Inches(0.85), Inches(0.95), Inches(11.6), Inches(1.1),
         [{"text": title, "size": 30, "bold": True, "color": SLATE}])
    if sub:
        text(slide, Inches(0.85), Inches(1.85), Inches(11.6), Inches(0.6),
             [{"text": sub, "size": 15, "color": MUTED}])


def refine_tag(slide):
    text(slide, Inches(11.0), Inches(6.95), Inches(2.0), Inches(0.4),
         [{"text": "[ I will refine ]", "size": 10, "italic": True, "color": CORAL, "align": PP_ALIGN.RIGHT}])


def footer(slide, n):
    text(slide, Inches(0.85), Inches(6.98), Inches(8), Inches(0.4),
         [{"text": "Verity: Payment Integrity Reasoning Copilot", "size": 9, "color": MUTED}])
    text(slide, Inches(12.2), Inches(6.98), Inches(0.9), Inches(0.4),
         [{"text": str(n), "size": 9, "color": MUTED, "align": PP_ALIGN.RIGHT}])


def bullets(slide, items, x=Inches(0.85), y=Inches(2.4), w=Inches(11.6), size=19, gap=12):
    runs = [{"text": t, "bullet": True, "size": size, "space_after": gap, "color": SLATE, "line_spacing": 1.05} for t in items]
    text(slide, x, y, w, Inches(4), runs)


# ── Slides ──────────────────────────────────────────────────────────────────
def s_title(prs):
    s = blank(prs)
    bg(s, PURPLE)
    box(s, 0, Inches(3.0), SW, Inches(0.06), fill=CORAL, radius=False)
    text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.4),
         [{"text": "Verity", "size": 54, "bold": True, "color": WHITE},
          {"text": "Payment Integrity Reasoning Copilot", "size": 24, "color": RGBColor(0xD9, 0xC7, 0xF0)}])
    text(s, Inches(0.9), Inches(3.4), Inches(11.5), Inches(1.6),
         [{"text": "An agentic AI that reviews a medical claim and returns an auditable PAY / FLAG / DENY recommendation, rules decide the facts, the model reasons and explains.",
           "size": 16, "color": RGBColor(0xEA, 0xE2, 0xF6), "line_spacing": 1.2}])
    text(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
         [{"text": "Srinidhi Jagannathan", "size": 16, "bold": True, "color": WHITE},
          {"text": "MS Business Analytics, Santa Clara University  ·  Built for Cotiviti", "size": 13, "color": RGBColor(0xC9, 0xB8, 0xE6)}])


def s_problem(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "THE PROBLEM")
    heading(s, "Payment-integrity review must be defensible, not just accurate")
    bullets(s, [
        "Denials get appealed, every decision needs a documented, cited rationale.",
        "Regulators expect explainability; black-box scoring is hard to defend.",
        "Wrong denials cause provider abrasion and rework.",
        "The bottleneck is explainable adjudication at scale.",
    ])
    refine_tag(s)
    footer(s, 2)


def s_topic(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "TOPIC FRAMING")
    heading(s, "Topic 2: Clinical Decision Making & Pattern Recognition")
    text(s, Inches(0.85), Inches(2.3), Inches(11.6), Inches(1.0),
         [{"text": "Verity targets the decision-making half for the Payment & Operations side of TPO:",
           "size": 17, "color": MUTED}])
    row = [("Chain reasoning", "step-by-step review a reviewer can follow"),
           ("Agentic generative AI", "an orchestrator that calls tools and decides"),
           ("Classification & inference", "claims sorted into PAY / FLAG / DENY")]
    x = Inches(0.85)
    for title, sub in row:
        c = box(s, x, Inches(3.3), Inches(3.7), Inches(2.0), fill=PURPLE_LT, line=LINE)
        c.line.fill.background()
        text(s, x + Inches(0.25), Inches(3.55), Inches(3.2), Inches(1.6),
             [{"text": title, "size": 17, "bold": True, "color": PURPLE},
              {"text": sub, "size": 13, "color": SLATE, "space_before": 4, "line_spacing": 1.1}])
        x = Emu(int(x) + int(Inches(3.9)))
    text(s, Inches(0.85), Inches(5.6), Inches(11.6), Inches(0.8),
         [{"text": "The statistical half (pattern recognition, clustering, time-series anomaly detection) is the complementary detection layer, addressed in the recommendation, not the POC.",
           "size": 13, "italic": True, "color": MUTED}])
    footer(s, 3)


def s_architecture(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "HOW VERITY WORKS")
    heading(s, "Rules decide facts. The model reasons and explains.")
    # Orchestrator
    orch = box(s, Inches(0.85), Inches(3.0), Inches(2.5), Inches(1.2), fill=PURPLE_LT, line=PURPLE, line_w=1.5)
    text(s, Inches(0.85), Inches(3.15), Inches(2.5), Inches(1.0),
         [{"text": "LLM Orchestrator", "size": 15, "bold": True, "color": PURPLE, "align": PP_ALIGN.CENTER},
          {"text": "reasons, calls tools", "size": 11, "color": MUTED, "align": PP_ALIGN.CENTER}])
    # 5 rule tools
    tools = ["Bundling", "Medical necessity", "Frequency / units", "Modifier", "Duplicate"]
    ty = Inches(1.5)
    for t in tools:
        b = box(s, Inches(4.9), ty, Inches(3.4), Inches(0.78), fill=WHITE, line=LINE, line_w=1.25)
        text(s, Inches(5.05), ty + Inches(0.16), Inches(3.1), Inches(0.5),
             [{"text": t, "size": 13, "bold": True, "color": SLATE}])
        ty = Emu(int(ty) + int(Inches(0.95)))
    # Verdict
    box(s, Inches(9.9), Inches(3.0), Inches(2.6), Inches(1.2), fill=RGBColor(0xF0, 0xFD, 0xF4), line=GREEN, line_w=1.5)
    text(s, Inches(9.9), Inches(3.15), Inches(2.6), Inches(1.0),
         [{"text": "Verdict", "size": 15, "bold": True, "color": GREEN, "align": PP_ALIGN.CENTER},
          {"text": "PAY / FLAG / DENY", "size": 11, "color": MUTED, "align": PP_ALIGN.CENTER}])
    # connectors
    conn1 = s.shapes.add_connector(2, Inches(3.35), Inches(3.6), Inches(4.9), Inches(3.0))
    conn1.line.color.rgb = LINE; conn1.line.width = Pt(1.5)
    conn2 = s.shapes.add_connector(2, Inches(8.3), Inches(3.0), Inches(9.9), Inches(3.6))
    conn2.line.color.rgb = LINE; conn2.line.width = Pt(1.5)
    text(s, Inches(0.85), Inches(5.7), Inches(11.6), Inches(1.0),
         [{"text": "Deterministic rules engine = the auditable backbone. The orchestrator decides which tools to call, interprets ambiguity, and synthesizes a cited rationale. Every step is recorded.",
           "size": 14, "color": MUTED, "line_spacing": 1.15}])
    footer(s, 4)


def s_rules(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "THE FIVE RULE CHECKS")
    heading(s, "Deterministic, cited, auditable")
    items = [
        ("Procedure-to-procedure (NCCI-style)", "flags bundled/mutually-exclusive CPT pairs billed together"),
        ("Medical necessity", "confirms the ICD-10 diagnosis supports the CPT procedure"),
        ("Frequency / units (MUE)", "flags units beyond the clinically plausible per-day maximum"),
        ("Modifier validation", "checks modifiers are recognized and required ones are present"),
        ("Duplicate detection", "catches the same member / procedure / date submitted twice"),
    ]
    y = Inches(2.45)
    for i, (t, d) in enumerate(items, 1):
        num = box(s, Inches(0.85), y, Inches(0.5), Inches(0.5), fill=PURPLE, radius=False)
        num.line.fill.background()
        text(s, Inches(0.85), y + Inches(0.04), Inches(0.5), Inches(0.5),
             [{"text": str(i), "size": 16, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
        text(s, Inches(1.55), y - Inches(0.02), Inches(11), Inches(0.7),
             [{"text": t, "size": 16, "bold": True, "color": SLATE},
              {"text": d, "size": 13, "color": MUTED}])
        y = Emu(int(y) + int(Inches(0.86)))
    footer(s, 5)


def s_nearmiss(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "THE STANDOUT MOMENT")
    heading(s, "The near-miss save, it avoids false denials, not just catches violations")
    # Fire box
    box(s, Inches(0.85), Inches(2.7), Inches(4.6), Inches(2.6), fill=RGBColor(0xFE, 0xF2, 0xF2), line=RED, line_w=1.5)
    text(s, Inches(1.1), Inches(2.95), Inches(4.1), Inches(2.2),
         [{"text": "1 · Edit fires", "size": 16, "bold": True, "color": RED},
          {"text": "Two same-day 93000 (EKG) lines.", "size": 14, "color": SLATE, "space_before": 6},
          {"text": "MUE per-line maximum is 1/day → on its face, an apparent overpayment.", "size": 14, "color": SLATE, "line_spacing": 1.15}])
    # arrow
    ar = s.shapes.add_shape(13, Inches(5.65), Inches(3.8), Inches(1.9), Inches(0.5))
    ar.fill.solid(); ar.fill.fore_color.rgb = PURPLE; ar.line.fill.background(); ar.shadow.inherit = False
    text(s, Inches(5.55), Inches(3.35), Inches(2.1), Inches(0.4),
         [{"text": "modifier 76", "size": 12, "bold": True, "color": PURPLE, "align": PP_ALIGN.CENTER}])
    # Clear box
    box(s, Inches(7.75), Inches(2.7), Inches(4.75), Inches(2.6), fill=RGBColor(0xF0, 0xFD, 0xF4), line=GREEN, line_w=1.5)
    text(s, Inches(8.0), Inches(2.95), Inches(4.25), Inches(2.2),
         [{"text": "2 · Override clears it → PAY", "size": 16, "bold": True, "color": GREEN},
          {"text": "Modifier 76 documents a justified repeat (same physician), billed on a separate line.", "size": 14, "color": SLATE, "space_before": 6, "line_spacing": 1.15},
          {"text": "Both the fire and the clear are recorded in the audit trail.", "size": 13, "italic": True, "color": MUTED, "space_before": 6}])
    text(s, Inches(0.85), Inches(5.75), Inches(11.6), Inches(0.8),
         [{"text": "False-positive avoidance and provider-abrasion sensitivity, the maturity that separates a copilot from a blunt auto-denier.",
           "size": 14, "bold": True, "color": PURPLE}])
    footer(s, 6)


def s_audit(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "AUDITABILITY")
    heading(s, "Every decision defends itself")
    bullets(s, [
        "Shows its full reasoning chain, step by step.",
        "Cites the specific rule and policy it applied.",
        "Exports a complete, ordered audit trail (JSON).",
        "The fire-then-clear is captured, the record proves the system reconsidered.",
    ])
    box(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(0.95), fill=PURPLE_LT, line=None)
    text(s, Inches(1.1), Inches(5.72), Inches(11.1), Inches(0.7),
         [{"text": "In a regulated payment context, the audit trail is the product, it is what survives an appeal.",
           "size": 15, "bold": True, "color": PURPLE}])
    footer(s, 7)


def s_eval(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "HOW IT'S EVALUATED")
    heading(s, "Measured against a labeled gold set")
    text(s, Inches(0.85), Inches(2.5), Inches(5.4), Inches(2.0),
         [{"text": "10 / 10", "size": 48, "bold": True, "color": PURPLE},
          {"text": "verdicts match the labeled ground truth on the sample set", "size": 15, "color": MUTED, "line_spacing": 1.15}])
    bullets(s, [
        "Each claim carries a ground-truth PAY/FLAG/DENY label.",
        "Per-rule breakdown shows which rule drove each verdict.",
        "Scales to precision / recall per disposition on a larger gold set.",
    ], x=Inches(6.6), y=Inches(2.6), w=Inches(6.0), size=16, gap=14)
    footer(s, 8)


def s_oppthreat(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "OPPORTUNITIES & THREATS")
    heading(s, "Why the auditable design is the point")
    box(s, Inches(0.85), Inches(2.5), Inches(5.6), Inches(3.6), fill=RGBColor(0xF0, 0xFD, 0xF4), line=GREEN)
    text(s, Inches(1.1), Inches(2.7), Inches(5.1), Inches(3.2),
         [{"text": "Opportunity", "size": 18, "bold": True, "color": GREEN},
          {"text": "Faster, defensible adjudication", "size": 14, "bullet": True, "color": SLATE, "space_before": 8},
          {"text": "Reduced provider abrasion via false-positive avoidance", "size": 14, "bullet": True, "color": SLATE},
          {"text": "Audit-ready rationale for appeals", "size": 14, "bullet": True, "color": SLATE}])
    box(s, Inches(6.9), Inches(2.5), Inches(5.6), Inches(3.6), fill=RGBColor(0xFE, 0xF2, 0xF2), line=RED)
    text(s, Inches(7.15), Inches(2.7), Inches(5.1), Inches(3.2),
         [{"text": "Threat", "size": 18, "bold": True, "color": RED},
          {"text": "Confident-but-unexplainable AI decisions", "size": 14, "bullet": True, "color": SLATE, "space_before": 8},
          {"text": "Hallucination if a model decides facts", "size": 14, "bullet": True, "color": SLATE},
          {"text": "PHI / compliance risk", "size": 14, "bullet": True, "color": SLATE},
          {"text": "Mitigated by: rules decide facts; model only explains; human-in-the-loop", "size": 13, "italic": True, "color": MUTED, "space_before": 8}])
    refine_tag(s)
    footer(s, 9)


def s_reco(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "STRATEGIC RECOMMENDATION")
    heading(s, "Agent-assisted, human-in-the-loop, auditable adjudication")
    bullets(s, [
        "AI drafts the reasoned, cited recommendation.",
        "A human reviewer confirms or overrides.",
        "The audit trail is the deliverable, defensible end to end.",
    ])
    box(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.4), fill=PURPLE, line=None)
    text(s, Inches(1.15), Inches(5.25), Inches(11.0), Inches(0.95),
         [{"text": "Not \u201creplace reviewers\u201d, make reviewers faster and decisions defensible.",
           "size": 20, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
    refine_tag(s)
    footer(s, 10)


def s_production(prs):
    s = blank(prs)
    bg(s, WHITE)
    kicker(s, "WHAT A PRODUCTION VERSION ADDS")
    heading(s, "From POC to platform")
    bullets(s, [
        "HIPAA / PHI handling, access control, and audit logging.",
        "A human review queue with accept / override and provider appeals.",
        "An evaluation pipeline and drift monitoring (models and policies change).",
    ], y=Inches(2.3), gap=10)
    text(s, Inches(0.85), Inches(4.55), Inches(11.6), Inches(0.5),
         [{"text": "Where the rest of Topic 2 fits, upstream of per-claim review:", "size": 15, "bold": True, "color": PURPLE}])
    bullets(s, [
        "Classification / risk tiering to prioritize claims.",
        "Clustering to surface suspicious provider cohorts.",
        "Time-series anomaly detection for billing aberrancy.",
    ], y=Inches(5.1), gap=8, size=15)
    footer(s, 11)


def s_close(prs):
    s = blank(prs)
    bg(s, PURPLE)
    text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.0),
         [{"text": "Rules decide the facts.", "size": 34, "bold": True, "color": WHITE},
          {"text": "The model reasons, explains, and leaves an audit trail.", "size": 34, "bold": True, "color": RGBColor(0xD9, 0xC7, 0xF0)}])
    text(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(0.8),
         [{"text": "Thank you., Srinidhi Jagannathan", "size": 18, "color": RGBColor(0xEA, 0xE2, 0xF6)}])


def main() -> None:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    for fn in [s_title, s_problem, s_topic, s_architecture, s_rules, s_nearmiss,
               s_audit, s_eval, s_oppthreat, s_reco, s_production, s_close]:
        fn(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
